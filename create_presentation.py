"""
Generate a styled PowerPoint presentation for the thesis project plan.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ──────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
ACCENT  = RGBColor(0x00, 0xD4, 0xAA)       # Teal/mint accent
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)       # Purple accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCC, 0xCC, 0xCC)
MUTED   = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0x24, 0x24, 0x3E)       # Slightly lighter card
WARM    = RGBColor(0xFF, 0x6B, 0x6B)        # Warm red for emphasis

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT, spacing=Pt(6), bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u25CF'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'[1:] if str(bullet_color).startswith('#') else format(bullet_color[0], '02X') + format(bullet_color[1], '02X') + format(bullet_color[2], '02X')})
        buClr.append(srgbClr)
        pPr.append(buClr)
        pPr.append(buChar)

    return txBox


def add_simple_bullets(slide, left, top, width, height, items, font_size=14,
                       color=LIGHT, spacing=Pt(4)):
    """Simpler bullet list using dash prefix."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing

    return txBox


# ── Create Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

blank_layout = prs.slide_layouts[6]  # Blank layout

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

# Accent bar at top
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

# Decorative circle (top right)
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT2
circle.line.fill.background()
circle.fill.fore_color.brightness = 0.7

# Small decorative circle
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(2.7), Inches(1), Inches(1))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT
circle2.line.fill.background()
circle2.fill.fore_color.brightness = 0.8

# Title
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(9), Inches(1.2),
             "HUMAN vs. AI", font_size=52, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.8), Inches(9), Inches(1),
             "Analysing the Advancement of AI in Music Mastering",
             font_size=28, color=ACCENT, bold=False)

# Divider line
add_accent_bar(slide, Inches(0.8), Inches(3.9), Inches(2), Inches(0.04), ACCENT)

# Subtitle info
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(0.5),
             "MSc Music and Media Technologies  |  Trinity College Dublin",
             font_size=16, color=LIGHT)

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(8), Inches(0.5),
             "Michael Adams", font_size=20, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(5.4), Inches(8), Inches(0.5),
             "Project Plan Presentation  |  2025", font_size=14, color=MUTED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Research Question & Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "RESEARCH QUESTION", font_size=14, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1.2),
             "Has AI mastering technology advanced to a point where it is\nperceptually indistinguishable from professional human mastering?",
             font_size=28, color=WHITE, bold=True)

add_accent_bar(slide, Inches(0.8), Inches(2.4), Inches(2), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.8),
             "This study conducts a rigorous blind comparison between professionally mastered\ntracks and AI-mastered equivalents, evaluated by diverse listener groups.",
             font_size=16, color=LIGHT)

# Three key cards
card_data = [
    ("Comparative", "Blind A/B testing of human\nvs. AI mastered tracks"),
    ("Perceptual", "Listener scoring across\nmultiple audio quality criteria"),
    ("Statistical", "Data-driven analysis with\ndiverse participant groups"),
]

card_width = Inches(3.5)
card_height = Inches(2.2)
start_x = Inches(0.8)
card_y = Inches(4.0)
gap = Inches(0.5)

for i, (title, desc) in enumerate(card_data):
    x = start_x + i * (card_width + gap)
    card = add_shape(slide, x, card_y, card_width, card_height, CARD_BG, 0.05)
    add_accent_bar(slide, x, card_y, card_width, Inches(0.04), ACCENT if i == 0 else ACCENT2 if i == 1 else WARM)
    add_text_box(slide, x + Inches(0.3), card_y + Inches(0.3), card_width - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), card_y + Inches(1.0), card_width - Inches(0.6), Inches(1),
                 desc, font_size=14, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Methodology Overview (5-Step Process)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "METHODOLOGY", font_size=14, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Five-Stage Research Process", font_size=32, color=WHITE, bold=True)

steps = [
    ("01", "Collect", "Curate diverse\npre-mastered songs", ACCENT),
    ("02", "Human Master", "Engage professional\nmastering engineers", ACCENT2),
    ("03", "AI Master", "Process through\nleading AI platforms", RGBColor(0xFF, 0xAA, 0x00)),
    ("04", "Blind Test", "Structured listening\ntests with scoring", WARM),
    ("05", "Analyse", "Statistical analysis\n& conclusions", RGBColor(0x00, 0x99, 0xFF)),
]

step_width = Inches(2.1)
step_height = Inches(3.2)
start_x = Inches(0.6)
step_y = Inches(2.2)
gap = Inches(0.25)

for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (step_width + gap)

    # Card
    card = add_shape(slide, x, step_y, step_width, step_height, CARD_BG, 0.05)

    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), step_y + Inches(0.3), Inches(0.8), Inches(0.8))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()

    # Number text
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    # Title
    add_text_box(slide, x + Inches(0.15), step_y + Inches(1.3), step_width - Inches(0.3), Inches(0.5),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.15), step_y + Inches(1.9), step_width - Inches(0.3), Inches(1),
                 desc, font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Arrow between steps (except last)
    if i < len(steps) - 1:
        arrow_x = x + step_width + Inches(0.02)
        add_text_box(slide, arrow_x, step_y + Inches(1.2), Inches(0.25), Inches(0.5),
                     "\u25B6", font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: 8-Week Timeline Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "PROJECT TIMELINE", font_size=14, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "8-Week Plan Overview", font_size=32, color=WHITE, bold=True)

weeks = [
    ("WK 1", "Research &\nSong Selection", ACCENT),
    ("WK 2", "Engage\nEngineers", ACCENT),
    ("WK 3", "AI Mastering\n& Test Design", ACCENT2),
    ("WK 4", "Normalise &\nTest Prep", ACCENT2),
    ("WK 5", "Begin Blind\nTesting", RGBColor(0xFF, 0xAA, 0x00)),
    ("WK 6", "Complete Tests\n& Early Analysis", RGBColor(0xFF, 0xAA, 0x00)),
    ("WK 7", "Deep Analysis\n& Writing", WARM),
    ("WK 8", "Conclusions &\nSubmission", WARM),
]

# Phase labels
phases = [
    (0, 2, "PREPARATION", ACCENT),
    (2, 2, "PRODUCTION", ACCENT2),
    (4, 2, "TESTING", RGBColor(0xFF, 0xAA, 0x00)),
    (6, 2, "ANALYSIS", WARM),
]

block_w = Inches(1.4)
block_h = Inches(2.0)
start_x = Inches(0.5)
block_y = Inches(3.0)
gap = Inches(0.12)

# Phase bars
phase_y = Inches(2.2)
for start_i, count, label, color in phases:
    px = start_x + start_i * (block_w + gap)
    pw = count * block_w + (count - 1) * gap
    bar = add_shape(slide, px, phase_y, pw, Inches(0.45), color, 0.1)
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# Week blocks
for i, (wk, desc, color) in enumerate(weeks):
    x = start_x + i * (block_w + gap)

    card = add_shape(slide, x, block_y, block_w, block_h, CARD_BG, 0.05)
    add_accent_bar(slide, x, block_y, block_w, Inches(0.04), color)

    add_text_box(slide, x + Inches(0.05), block_y + Inches(0.2), block_w - Inches(0.1), Inches(0.4),
                 wk, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.05), block_y + Inches(0.7), block_w - Inches(0.1), Inches(1.2),
                 desc, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Key milestones
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
             "KEY MILESTONES", font_size=12, color=ACCENT, bold=True)

milestones = [
    "Week 2: Engineers engaged & AI platforms selected",
    "Week 4: All masters collected, normalised & test platform ready",
    "Week 6: All listening tests completed",
    "Week 8: Thesis submitted",
]

for i, m in enumerate(milestones):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.5)
    y = Inches(5.95) + row * Inches(0.4)
    add_text_box(slide, x, y, Inches(5.5), Inches(0.35),
                 f"\u2713  {m}", font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Weeks 1 & 2 Detail
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "WEEKS 1-2  |  PREPARATION", font_size=14, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Research, Song Selection & Engineer Engagement", font_size=28, color=WHITE, bold=True)

# Week 1 card
card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.04), ACCENT)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 1: Research & Song Selection", font_size=18, color=WHITE, bold=True)

w1_items = [
    "Survey academic literature on AI mastering",
    "Review prior human vs. AI audio studies",
    "Select 4-6 pre-mastered songs across genres",
    "Ensure genre diversity (rock, electronic, jazz, pop, etc.)",
    "Source high-quality unmastered mixes (24-bit WAV)",
    "Secure permissions for all tracks",
    "Submit ethics application for listening tests",
    "Prepare participant consent forms",
]
add_simple_bullets(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(3.5),
                   w1_items, font_size=12, color=LIGHT)

# Week 2 card
card = add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.04), ACCENT)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 2: Engage Mastering Engineers", font_size=18, color=WHITE, bold=True)

w2_items = [
    "Shortlist 2-3 professional mastering engineers",
    "Contact engineers with academic project brief",
    "Agree timelines, fees & deliverable formats",
    "Prepare standardised mastering briefs",
    "Provide identical source files to all engineers",
    "Request 16-bit/44.1kHz + 24-bit/48kHz masters",
    "Identify 2-3 AI mastering platforms to test",
    "Document AI platform features & capabilities",
]
add_simple_bullets(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
                   w2_items, font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Weeks 3 & 4 Detail
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT2)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "WEEKS 3-4  |  PRODUCTION", font_size=14, color=ACCENT2, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "AI Mastering, Normalisation & Test Preparation", font_size=28, color=WHITE, bold=True)

# Week 3 card
card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.04), ACCENT2)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 3: AI Mastering & Test Design", font_size=18, color=WHITE, bold=True)

w3_items = [
    "Process all songs through each AI platform",
    "Experiment with settings (presets, loudness, tone)",
    "Document all AI settings with screenshots",
    "Record processing times & limitations",
    "Export AI masters in matching formats",
    "Design listening test methodology (ABX/MUSHRA)",
    "Define scoring criteria (clarity, punch, warmth, etc.)",
    "Plan Likert scale (1-10) evaluation framework",
]
add_simple_bullets(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(3.5),
                   w3_items, font_size=12, color=LIGHT)

# Week 4 card
card = add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.04), ACCENT2)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 4: Normalise & Build Test Platform", font_size=18, color=WHITE, bold=True)

w4_items = [
    "Collect all completed human masters",
    "Loudness-normalise all tracks (e.g., -14 LUFS)",
    "Trim tracks to identical start/end points",
    "Assign anonymous labels (Version A, B, C...)",
    "Set up listening test platform (webMUSHRA / Forms)",
    "Upload all anonymised audio files",
    "Create questionnaire with demographic questions",
    "Run pilot test with 2-3 people & refine",
]
add_simple_bullets(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
                   w4_items, font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Weeks 5 & 6 Detail
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), RGBColor(0xFF, 0xAA, 0x00))

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "WEEKS 5-6  |  TESTING", font_size=14, color=RGBColor(0xFF, 0xAA, 0x00), bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Blind Listening Tests & Early Analysis", font_size=28, color=WHITE, bold=True)

# Week 5 card
card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.04), RGBColor(0xFF, 0xAA, 0x00))
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 5: Recruit & Begin Testing", font_size=18, color=WHITE, bold=True)

w5_items = [
    "Recruit from MMT classmates (trained ears)",
    "Recruit from TCD music students",
    "Engage professional musicians/engineers",
    "Include general listeners (non-musicians)",
    "Target 20-30 participants minimum",
    "Distribute & collect consent forms",
    "Conduct listening test sessions",
    "Keep sessions under 30-40 minutes",
]
add_simple_bullets(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(3.5),
                   w5_items, font_size=12, color=LIGHT)

# Week 6 card
card = add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.04), RGBColor(0xFF, 0xAA, 0x00))
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 6: Complete Tests & Early Analysis", font_size=18, color=WHITE, bold=True)

w6_items = [
    "Follow up with outstanding participants",
    "Close test window & finalise data collection",
    "Clean & validate all collected data",
    "Remove incomplete / flag outlier responses",
    "Calculate mean scores & standard deviations",
    "Compare human vs. AI across all criteria",
    "Break down by participant group",
    "Generate initial charts & visualisations",
]
add_simple_bullets(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
                   w6_items, font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Weeks 7 & 8 Detail
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), WARM)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "WEEKS 7-8  |  ANALYSIS & COMPLETION", font_size=14, color=WARM, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Deep Analysis, Conclusions & Submission", font_size=28, color=WHITE, bold=True)

# Week 7 card
card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.6), Inches(0.04), WARM)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 7: Deep Analysis & Writing", font_size=18, color=WHITE, bold=True)

w7_items = [
    "Overall preference analysis: Human vs. AI",
    "Per-criterion breakdown (clarity, dynamics, etc.)",
    "Genre-dependent performance analysis",
    "Expert vs. non-expert perception differences",
    "Platform-specific AI performance comparison",
    "Statistical significance testing (t-test, ANOVA)",
    "Optional: Spectral analysis (frequency, stereo width)",
    "Begin writing: Intro, Literature Review, Methodology",
]
add_simple_bullets(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(3.5),
                   w7_items, font_size=12, color=LIGHT)

# Week 8 card
card = add_shape(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.04), WARM)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5),
             "Week 8: Conclusions & Submission", font_size=18, color=WHITE, bold=True)

w8_items = [
    "Synthesise findings into clear conclusions",
    "Has AI mastering reached professional quality?",
    "Where does AI excel? Where does it fall short?",
    "Implications for the music industry",
    "Discuss limitations & future directions",
    "Complete Discussion & Conclusions chapters",
    "Write Abstract, compile references",
    "Final proofread, format & submit",
]
add_simple_bullets(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
                   w8_items, font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Tools & Technology
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "TOOLS & TECHNOLOGY", font_size=14, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Platforms, Software & Analysis Tools", font_size=28, color=WHITE, bold=True)

# AI Platforms card
card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(3.6), Inches(0.04), ACCENT)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(3), Inches(0.5),
             "AI Mastering Platforms", font_size=16, color=WHITE, bold=True)
ai_items = [
    "LANDR",
    "CloudBounce",
    "eMastered",
    "iZotope Ozone AI",
    "BandLab Mastering",
    "Dolby.io Media Mastering",
]
add_simple_bullets(slide, Inches(1.1), Inches(2.8), Inches(3), Inches(3.5),
                   ai_items, font_size=13, color=LIGHT)

# Analysis Tools card
card = add_shape(slide, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(4.8), Inches(2.0), Inches(3.6), Inches(0.04), ACCENT2)
add_text_box(slide, Inches(5.1), Inches(2.2), Inches(3), Inches(0.5),
             "Analysis Tools", font_size=16, color=WHITE, bold=True)
analysis_items = [
    "Python (librosa, scipy)",
    "matplotlib / seaborn",
    "SPSS or R for statistics",
    "iZotope Insight",
    "Voxengo SPAN",
    "Excel / Google Sheets",
]
add_simple_bullets(slide, Inches(5.1), Inches(2.8), Inches(3), Inches(3.5),
                   analysis_items, font_size=13, color=LIGHT)

# Test Platform card
card = add_shape(slide, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.6), CARD_BG, 0.03)
add_accent_bar(slide, Inches(8.8), Inches(2.0), Inches(3.6), Inches(0.04), RGBColor(0xFF, 0xAA, 0x00))
add_text_box(slide, Inches(9.1), Inches(2.2), Inches(3), Inches(0.5),
             "Listening Test Platforms", font_size=16, color=WHITE, bold=True)
test_items = [
    "webMUSHRA (open-source)",
    "Google Forms + audio",
    "Custom web application",
    "Quality headphones / monitors",
    "Treated listening room",
    "Reference DAW (Logic/Reaper)",
]
add_simple_bullets(slide, Inches(9.1), Inches(2.8), Inches(3), Inches(3.5),
                   test_items, font_size=13, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Risk Mitigation
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), WARM)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "RISK MITIGATION", font_size=14, color=WARM, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.6),
             "Identified Risks & Contingency Plans", font_size=28, color=WHITE, bold=True)

risks = [
    ("Engineer Delays", "Mastering engineers may not deliver on time",
     "Engage engineers early in Week 2; identify backup engineers", WARM),
    ("Low Participant Count", "Insufficient listeners weakens statistical significance",
     "Start recruitment in Week 3; use social media & course networks", RGBColor(0xFF, 0xAA, 0x00)),
    ("AI Platform Limitations", "Tools may not support certain genres or formats",
     "Test platforms early in Week 2; have backup platforms ready", ACCENT2),
    ("Test Platform Issues", "Technical problems during listening sessions",
     "Pilot test in Week 4; prepare offline backup (USB + paper)", ACCENT),
    ("Ethics Approval Delays", "Late approval blocks participant testing",
     "Submit ethics application in Week 1; follow up proactively", RGBColor(0x00, 0x99, 0xFF)),
]

card_h = Inches(0.85)
start_y = Inches(2.0)
gap = Inches(0.15)

for i, (risk, impact, mitigation, color) in enumerate(risks):
    y = start_y + i * (card_h + gap)

    card = add_shape(slide, Inches(0.8), y, Inches(11.5), card_h, CARD_BG, 0.02)

    # Color indicator bar on left
    add_accent_bar(slide, Inches(0.8), y, Inches(0.06), card_h, color)

    # Risk name
    add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 risk, font_size=14, color=WHITE, bold=True)

    # Impact
    add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(3.5), Inches(0.35),
                 impact, font_size=11, color=MUTED)

    # Mitigation
    add_text_box(slide, Inches(5.5), y + Inches(0.15), Inches(0.8), Inches(0.35),
                 "Mitigation:", font_size=10, color=color, bold=True)
    add_text_box(slide, Inches(6.4), y + Inches(0.15), Inches(5.5), Inches(0.55),
                 mitigation, font_size=12, color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11: Thank You / Questions
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), ACCENT)

# Decorative circles
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(5), Inches(3), Inches(3))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT2
circle.line.fill.background()
circle.fill.fore_color.brightness = 0.85

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(0.5), Inches(3.5), Inches(3.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT
circle2.line.fill.background()
circle2.fill.fore_color.brightness = 0.85

add_text_box(slide, Inches(2), Inches(2.2), Inches(9), Inches(1.2),
             "Thank You", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), ACCENT)

add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
             "Questions & Discussion", font_size=24, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.5),
             "Michael Adams  |  MSc Music & Media Technologies  |  Trinity College Dublin",
             font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
output_path = "/Users/michaeladams/Work/Esker/Development/AI/Claude/thesis-plan/Thesis_Project_Plan.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
